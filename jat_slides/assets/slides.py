import datetime
from upath import UPath as Path


import numpy as np
import pandas as pd
from babel.dates import format_date
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import SlideLayout
from pptx.text.text import TextFrame
from pptx.util import Cm, Pt

import dagster as dg
from dagster import AssetIn, asset
from jat_slides.partitions import zone_partitions
from jat_slides.resources import ConfigResource

from io import BytesIO
import fsspec
from cfc_core_utils import storage_options
from utils.cloud_paths import cloud_exists

RGB_BLUE = RGBColor(0x00, 0x70, 0xC0)
RGB_RED = RGBColor(0xFF, 0x00, 0x00)

def insert_picture_any(figure_shape, picture_path):
    # local file → use path
    if getattr(picture_path, "protocol", "file") == "file":
        return figure_shape.insert_picture(str(picture_path))
    # az:// → stream bytes via fsspec
    with fsspec.open(str(picture_path), "rb", **storage_options(picture_path)) as f:
        buf = BytesIO(f.read())
    return figure_shape.insert_picture(buf)

def find_layouts(pres) -> dict:
    layouts = {}
    for layout in pres.slide_layouts:
        name = layout.name
        if name == "Title Slide":
            layouts["title"] = layout
        elif name == "Picture with Title":
            layouts["picture_with_title"] = layout
        elif name == "Picture with Title and Content":
            layouts["picture_with_title_and_content"] = layout
        elif name == "Picture with Title, Content and Table":
            layouts["picture_with_title_content_table"] = layout
        elif name == "Section Divider":
            layouts["section"] = layout
    return layouts


def find_shape(shapes: SlideShapes, prefix: str) -> BaseShape:
    for shape in shapes:
        if shape.name.startswith(prefix):
            return shape

    raise ValueError(
        f"Shape with prefix '{prefix}' not found in the slide shapes.",
    )


def add_title_slide(pres: PresentationType, title_layout: SlideLayout) -> None:
    title = pres.slides.add_slide(title_layout)

    placeholders = title.shapes
    title_shape = placeholders[0]
    date_shape = placeholders[2]

    if not isinstance(title_shape, Shape) or not isinstance(date_shape, Shape):
        raise TypeError("Placeholders are not of type Shape.")

    title_shape.text = "Dinámicas Urbanas"
    date_shape.text = format_date(
        datetime.date.today(),
        locale="es",
        format="long",
    )


def add_section_slide(
    pres: PresentationType,
    section_layout: SlideLayout,
    section_name: str,
) -> None:
    section = pres.slides.add_slide(section_layout)
    section_title_shape = section.shapes[0]

    if not isinstance(section_title_shape, Shape):
        err = "Section title shape is not of type Shape."
        raise TypeError(err)
    section_title_shape.text = section_name


def add_highlighted_text_to_frame(
    shape: Shape,
    *,
    text_left: str,
    highlight: str,
    text_right: str,
    color: RGBColor,
) -> None:
    frame: TextFrame = shape.text_frame
    frame.clear()

    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text_left

    run = p.add_run()
    run.text = highlight
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = color

    run = p.add_run()
    run.text = text_right


def add_normal_text_to_shape(shape: Shape, text: str) -> None:
    frame: TextFrame = shape.text_frame
    frame.clear()

    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text


def add_picture_with_highlight_slide(
    pres: PresentationType,
    layout: SlideLayout,
    *,
    title: str,
    text_left: str,
    text_right: str,
    highlight: str,
    picture_path: Path,
    color: RGBColor,
) -> None:
    pop_slide = pres.slides.add_slide(layout)

    title_shape = find_shape(pop_slide.shapes, prefix="Text Placeholder 2")
    add_normal_text_to_shape(title_shape, title)

    text_shape = find_shape(pop_slide.shapes, prefix="Text Placeholder 3")
    add_highlighted_text_to_frame(
        text_shape,
        text_left=text_left,
        highlight=highlight,
        text_right=text_right,
        color=color,
    )

    figure_shape = find_shape(pop_slide.shapes, prefix="Picture")
    insert_picture_any(figure_shape, picture_path)


# pylint: disable=protected-access
def add_built_slide(
    pres: PresentationType,
    layout: SlideLayout,
    built_after_frac: float,
    *,
    pop_df: pd.DataFrame,
    built_area_df: pd.DataFrame,
    urban_area_df: pd.DataFrame,
    picture_path: Path,
) -> None:
    pop_df = pop_df.set_index("year")["pop"] / 1e6
    built_area_df = built_area_df.set_index("year")["area"] / 1e6
    urban_area_df = urban_area_df.set_index("year")["area"] / 1e6

    pop_change = (pop_df.loc[2020] - pop_df.loc[2000]) / pop_df.loc[2000] + 1
    built_change = (
        built_area_df.loc[2020] - built_area_df.loc[2000]
    ) / built_area_df.loc[2000] + 1
    urban_change = (
        urban_area_df.loc[2020] - urban_area_df.loc[2000]
    ) / urban_area_df.loc[2000] + 1

    built_slide = pres.slides.add_slide(layout)

    title_shape = find_shape(built_slide.shapes, prefix="Text Placeholder 2")
    add_normal_text_to_shape(title_shape, "Dinámica de urbanización 2000 a 2020")

    text_shape = find_shape(built_slide.shapes, prefix="Text Placeholder 3")
    add_highlighted_text_to_frame(
        text_shape,
        text_left="El ",
        highlight=f"{built_after_frac:.0%}",
        text_right=" de los píxeles se urbanizaron en el año 2000 o posterior.",
        color=RGB_BLUE,
    )

    table_shape = find_shape(built_slide.shapes, prefix="Table")
    table_frame = table_shape.insert_table(rows=4, cols=4)
    table = table_frame.table

    table.columns[0].width = Cm(1.8)
    table.columns[1].width = Cm(3.0)
    table.columns[2].width = Cm(3.0)
    table.columns[3].width = Cm(2.9)

    tbl = table_frame._element.graphic.graphicData.tbl
    tbl[0][-1].text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

    text_arr = np.array(
        [
            [
                "",
                "Millones de habitantes",
                "Superficie  construida (km²)",
                "Superficie urbana (km²)",
            ],
            [
                "2000",
                f"{pop_df.loc[2000]:.2f}",
                f"{built_area_df.loc[2000]:.1f}",
                f"{urban_area_df.loc[2000]:.1f}",
            ],
            [
                "2020",
                f"{pop_df.loc[2020]:.2f}",
                f"{built_area_df.loc[2020]:.1f}",
                f"{urban_area_df.loc[2020]:.1f}",
            ],
            ["", f"x{pop_change:.2f}", f"x{built_change:.2f}", f"x{urban_change:.2f}"],
        ],
    )

    style_arr = np.array(
        [
            [False, True, True, True],
            [True, False, False, False],
            [True, False, False, False],
            [False, True, True, True],
        ],
    )

    for row_idx in range(text_arr.shape[0]):
        for col_idx in range(text_arr.shape[1]):
            cell = table.cell(row_idx, col_idx)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text_arr[row_idx, col_idx]
            run.font.size = Pt(14)

            if style_arr[row_idx, col_idx]:
                run.font.bold = True
                run.font.color.rgb = RGB_BLUE

    if cloud_exists(picture_path):
        figure_shape = find_shape(built_slide.shapes, prefix="Picture")
        insert_picture_any(figure_shape, picture_path)


def add_single_picture_slide(
    pres: PresentationType,
    layout: SlideLayout,
    *,
    picture_path: Path,
    title: str,
    total_jobs: float | None = None,
) -> None:
    slide = pres.slides.add_slide(layout)

    title_shape = find_shape(slide.shapes, prefix="Text")
    add_normal_text_to_shape(title_shape, title)

    if cloud_exists(picture_path):
        figure_shape = find_shape(slide.shapes, prefix="Picture")
        insert_picture_any(figure_shape, picture_path)


def generate_single_slide(
    *,
    name: str,
    lost_pop_after_2000: float,
    built_after_2000: float,
    total_jobs: float,
    pop_df: pd.DataFrame,
    built_df: pd.DataFrame,
    built_urban_df: pd.DataFrame,
    built_figure_path: Path,
    pg_figure_path: Path,
    income_figure_path: Path | None,
    jobs_figure_path: Path,
) -> PresentationType:
    pres = Presentation("./template.pptx")
    layouts = find_layouts(pres)

    add_section_slide(pres, layouts["section"], name)

    add_picture_with_highlight_slide(
        pres,
        layouts["picture_with_title_and_content"],
        title="Dinámica de población 2000 a 2020",
        text_left="El ",
        highlight=f"{lost_pop_after_2000:.0%}",
        text_right=" de la superficie construida de 2020 perdió población con respecto al 2000.",
        picture_path=pg_figure_path,
        color=RGB_RED,
    )

    add_built_slide(
        pres,
        layouts["picture_with_title_content_table"],
        built_after_2000,
        pop_df=pop_df,
        built_area_df=built_df,
        urban_area_df=built_urban_df,
        picture_path=built_figure_path,
    )

    if income_figure_path is not None:
        add_single_picture_slide(
            pres,
            layouts["picture_with_title"],
            picture_path=income_figure_path,
            title="Ingreso",
        )

    add_picture_with_highlight_slide(
        pres,
        layouts["picture_with_title_and_content"],
        title="Empleos totales",
        text_left="La zona metropolitana tiene un total de ",
        highlight=f"{total_jobs:,.0f}",
        text_right=" empleos.",
        picture_path=jobs_figure_path,
        color=RGB_BLUE,
    )

    # add_single_picture_slide(pres, layouts["picture_with_title_and_content"], picture_path=jobs_figure_path, title="Número total de empleos")

    return pres


@asset(
    ins={
        "lost_pop_after_2000": AssetIn(key=["stats", "lost_pop_after_2000"]),
        "built_after_2000": AssetIn(key=["stats", "built_after_2000"]),
        "pop_df": AssetIn(key=["stats", "population"]),
        "total_jobs": AssetIn(key=["stats", "total_jobs"]),
        "built_df": AssetIn(key=["stats", "built_area"]),
        "built_urban_df": AssetIn(key=["stats", "built_urban_area"]),
        "pg_figure_path": AssetIn(
            key=["plot", "population_grid"],
            input_manager_key="path_manager",
        ),
        "built_figure_path": AssetIn(
            key=["plot", "built"],
            input_manager_key="path_manager",
        ),
        "income_figure_path": AssetIn(
            key=["plot", "income"],
            input_manager_key="path_manager",
        ),
        "jobs_figure_path": AssetIn(
            key=["plot", "jobs"],
            input_manager_key="path_manager",
        ),
    },
    partitions_def=zone_partitions,
    io_manager_key="presentation_manager",
)
def slides(
    context: dg.AssetExecutionContext,
    zone_config_resource: ConfigResource,
    lost_pop_after_2000: float,
    built_after_2000: float,
    total_jobs: float,
    pop_df: pd.DataFrame,
    built_df: pd.DataFrame,
    built_urban_df: pd.DataFrame,
    pg_figure_path: Path,
    built_figure_path: Path,
    income_figure_path: Path,
    jobs_figure_path: Path,
) -> PresentationType:
    return generate_single_slide(
        name=zone_config_resource.names[context.partition_key],
        lost_pop_after_2000=lost_pop_after_2000,
        built_after_2000=built_after_2000,
        total_jobs=total_jobs,
        pop_df=pop_df,
        built_df=built_df,
        built_urban_df=built_urban_df,
        built_figure_path=built_figure_path,
        pg_figure_path=pg_figure_path,
        income_figure_path=income_figure_path,
        jobs_figure_path=jobs_figure_path,
    )
