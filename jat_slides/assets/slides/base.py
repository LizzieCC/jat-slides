import datetime

import numpy as np
import pandas as pd

from babel.dates import format_date
from jat_slides.resources import PathResource, ZonesResource
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import SlideLayout
from pptx.text.text import TextFrame
from pptx.util import Cm, Pt
from dagster import asset, AssetIn
from pathlib import Path


def find_layouts(pres: PresentationType) -> dict[str, SlideLayout]:
    layouts = {}
    for layout in pres.slide_layouts:
        name = layout.name
        if name == "Title Slide":
            layouts["title"] = layout
        elif name == "Map and Content":
            layouts["pop"] = layout
        elif name == "Map and Table":
            layouts["built"] = layout
        elif name == "Section Divider":
            layouts["section"] = layout
    return layouts


def find_shape(shapes: SlideShapes, prefix: str):
    for shape in shapes:
        if shape.name.startswith(prefix):
            return shape


def add_title_slide(pres: PresentationType, title_layout: SlideLayout) -> None:
    title = pres.slides.add_slide(title_layout)

    placeholders = title.shapes
    placeholders[0].text = "Dinámicas Urbanas"
    placeholders[2].text = format_date(
        datetime.date.today(), locale="es", format="long"
    )


def add_section_slide(
    pres: PresentationType, section_layout: SlideLayout, section_name: str
) -> None:
    section = pres.slides.add_slide(section_layout)
    section.shapes[0].text = section_name


def add_highlighted_text_to_frame(
    frame: TextFrame, *, text_left: str, highlight: float, text_right: str
) -> None:
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text_left

    run = p.add_run()
    run.text = f"{highlight:.0%}"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)

    run = p.add_run()
    run.text = text_right


def add_lost_pop_slide(
    pres: PresentationType,
    layout: SlideLayout,
    lost_pop_frac: float,
    picture_path: Path,
) -> None:
    pop_slide = pres.slides.add_slide(layout)

    text_shape = find_shape(pop_slide.shapes, prefix="Text")
    frame: TextFrame = text_shape.text_frame
    frame.clear()

    add_highlighted_text_to_frame(
        frame,
        text_left="El ",
        highlight=lost_pop_frac,
        text_right=" de la superficie de 2020 perdió población con respecto al 2000.",
    )

    if picture_path.exists():
        figure_shape = find_shape(pop_slide.shapes, prefix="Picture")
        figure_shape.insert_picture(str(picture_path))


# pylint: disable=protected-access
def add_built_slide(
    pres: PresentationType,
    layout: SlideLayout,
    built_after_frac: float,
    *,
    pop_df: pd.DataFrame,
    built_area_df: pd.DataFrame,
    urban_area_df: pd.DataFrame,
    picture_path: Path,
) -> None:
    pop_df = pop_df.set_index("year")["pop"] / 1e6
    built_area_df = built_area_df.set_index("year")["area"] / 1e6
    urban_area_df = urban_area_df.set_index("year")["area"] / 1e6

    pop_change = (pop_df.loc[2020] - pop_df.loc[2000]) / pop_df.loc[2000] + 1
    built_change = (
        built_area_df.loc[2020] - built_area_df.loc[2000]
    ) / built_area_df.loc[2000] + 1
    urban_change = (
        urban_area_df.loc[2020] - urban_area_df.loc[2000]
    ) / urban_area_df.loc[2000] + 1

    built_slide = pres.slides.add_slide(layout)

    text_shape = find_shape(built_slide.shapes, prefix="Text")
    frame: TextFrame = text_shape.text_frame
    frame.clear()

    add_highlighted_text_to_frame(
        frame,
        text_left="El ",
        highlight=built_after_frac,
        text_right=" de los píxeles se urbanizaron en el año 2000 o posterior.",
    )

    table_shape = find_shape(built_slide.shapes, prefix="Table")
    table_frame = table_shape.insert_table(rows=4, cols=4)
    table = table_frame.table

    table.columns[0].width = Cm(1.8)
    table.columns[1].width = Cm(3.0)
    table.columns[2].width = Cm(3.0)
    table.columns[3].width = Cm(2.9)

    tbl = table_frame._element.graphic.graphicData.tbl
    tbl[0][-1].text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

    text_arr = np.array(
        [
            [
                "",
                "Millones de habitantes",
                "Superficie  construida (km²)",
                "Superficie urbana (km²)",
            ],
            [
                "2000",
                f"{pop_df.loc[2000]:.2f}",
                f"{built_area_df.loc[2000]:.1f}",
                f"{urban_area_df.loc[2000]:.1f}",
            ],
            [
                "2020",
                f"{pop_df.loc[2020]:.2f}",
                f"{built_area_df.loc[2020]:.1f}",
                f"{urban_area_df.loc[2020]:.1f}",
            ],
            ["", f"x{pop_change:.2f}", f"x{built_change:.2f}", f"x{urban_change:.2f}"],
        ]
    )

    style_arr = np.array(
        [
            [False, True, True, True],
            [True, False, False, False],
            [True, False, False, False],
            [False, True, True, True],
        ]
    )

    for row_idx in range(text_arr.shape[0]):
        for col_idx in range(text_arr.shape[1]):
            cell = table.cell(row_idx, col_idx)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text_arr[row_idx, col_idx]
            run.font.size = Pt(14)

            if style_arr[row_idx, col_idx]:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)

    if picture_path.exists():
        figure_shape = find_shape(built_slide.shapes, prefix="Picture")
        figure_shape.insert_picture(str(picture_path))


@asset(
    ins={
        "lost_pop_after_2000": AssetIn(key=["stats", "lost_pop_after_2000"]),
        "built_after_2000": AssetIn(key=["stats", "built_after_2000"]),
        "pop_df": AssetIn(key=["stats", "population"]),
        "built_df": AssetIn(key=["stats", "built_area"]),
        "built_urban_df": AssetIn(key=["stats", "built_urban_area"]),
    },
    io_manager_key="presentation_manager",
)
def slides(
    path_resource: PathResource,
    zones_resource: ZonesResource,
    lost_pop_after_2000: pd.DataFrame,
    built_after_2000: pd.DataFrame,
    pop_df: dict[str, pd.DataFrame],
    built_df: dict[str, pd.DataFrame],
    built_urban_df: dict[str, pd.DataFrame],
) -> PresentationType:
    figure_path = Path(path_resource.figure_path)

    pres = Presentation("./template.pptx")
    layouts = find_layouts(pres)

    add_title_slide(pres, layouts["title"])

    for zone in zones_resource.wanted_zones:
        add_section_slide(pres, layouts["section"], zones_resource.zone_names[zone])

        lost_pop_frac = lost_pop_after_2000.loc[
            lost_pop_after_2000["zone"] == zone, "lost"
        ].item()
        lost_pop_fig = figure_path / f"mesh/{zone}.png"
        add_lost_pop_slide(pres, layouts["pop"], lost_pop_frac, lost_pop_fig)

        built_after_frac = built_after_2000.loc[
            built_after_2000["zone"] == zone, "built"
        ].item()
        built_year_fig = figure_path / f"built/{zone}_r.png"
        add_built_slide(
            pres,
            layouts["built"],
            built_after_frac,
            pop_df=pop_df[zone],
            built_area_df=built_df[zone],
            urban_area_df=built_urban_df[zone],
            picture_path=built_year_fig,
        )

    return pres
